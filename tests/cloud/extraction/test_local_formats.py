# test_local_formats.py — coverage for the file types LocalExtractor gained on
# 2026-08-29 (files-intelligence Track 1): .pptx, .xlsx, .html and the raster
# set .webp/.tif/.tiff/.heic/.heif.
#
# Created 2026-08-29. Sibling of test_local.py, which pins the original
# PDF/DOCX/image/VTT parity contract and is left untouched.
#
# THREE THINGS THESE TESTS DO ON PURPOSE:
#
# 1. Real fixtures, built at test time. Every document here is written by the
#    library that will read it back (python-pptx, openpyxl, Pillow), the same
#    way test_local.py builds its PDF with PdfWriter. A mocked parser proves
#    the branch was entered; it does not prove the branch works.
# 2. Only pytesseract is mocked. The tesseract BINARY is not guaranteed on any
#    machine, so the OCR call is stubbed — but PIL decodes the real webp/tiff/
#    heic bytes, and the assertions check the format of the image the stub
#    received. That is what proves the fixture reached the branch rather than
#    a mock answering for it.
# 3. Missing-dependency cases assert a LOUD RuntimeError. This module has
#    already shipped the silent version of that bug once (pypdf absent from the
#    production image: green build, every PDF empty, read as "the feature is
#    off"), so each new lazy import gets the same explicit test.
from __future__ import annotations

import sys
from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest
from pocketpaw_ee.cloud.extraction import local as local_mod
from pocketpaw_ee.cloud.extraction.local import LocalExtractor

PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"

# An article-shaped document. trafilatura is tuned for articles and returns
# None on thin or boilerplate-only markup, so a too-simple fixture would
# silently exercise the FALLBACK path while appearing to test trafilatura.
ARTICLE_HTML = """<html><head><title>Why extraction matters</title></head><body>
<nav>home about contact login</nav>
<article>
<h1>Why extraction matters</h1>
<p>A document that cannot be read is a document that cannot be searched.
The pipeline turns bytes into sentences and sentences into recall.</p>
<p>Every format the extractor skips is a silent hole in the knowledge base,
and a silent hole reads to a user as a product that simply does not work.</p>
<p>Coverage is therefore not a nice to have. It is the whole feature, and the
absence of coverage is indistinguishable from an outage.</p>
</article>
<footer>copyright 2026 all rights reserved</footer></body></html>"""

# Bodyless fragments — the shape trafilatura declines (returns None), which is
# the ONLY way to reach the tag-strip fallback. Verified by the guard test
# below; see its docstring for why the obvious "thin document" fixtures do not
# work for this.
FRAGMENT_HTML = "<div><b>visible</b> fragment text</div>"
FRAGMENT_WITH_SCRIPT_HTML = (
    "<div><script>var leak = 1;</script>"
    "<style>.a{color:crimson}</style><b>visible</b> fragment text</div>"
)


# --------------------------------------------------------------------------
# fixtures
# --------------------------------------------------------------------------
@pytest.fixture
def deck(tmp_path: Path) -> Path:
    """A real two-slide .pptx with body text, a table and speaker notes."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    first = prs.slides.add_slide(prs.slide_layouts[5])
    first.shapes.title.text = "Quarterly Revenue"
    box = first.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    box.text_frame.text = "Bookings grew eighteen percent"
    first.notes_slide.notes_text_frame.text = "Mention the Frankfurt outage"

    second = prs.slides.add_slide(prs.slide_layouts[5])
    second.shapes.title.text = "Pipeline"
    table = second.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(4), Inches(1)).table
    table.cell(0, 0).text = "Account"
    table.cell(0, 1).text = "Stage"
    table.cell(1, 0).text = "Acme"
    table.cell(1, 1).text = "Won"

    path = tmp_path / "deck.pptx"
    prs.save(str(path))
    return path


@pytest.fixture
def workbook(tmp_path: Path) -> Path:
    """A real two-sheet .xlsx."""
    from openpyxl import Workbook

    wb = Workbook()
    sheet = wb.active
    sheet.title = "Pipeline"
    sheet.append(["Account", "Stage", "ARR"])
    sheet.append(["Acme", "Won", 42000])
    notes = wb.create_sheet("Notes")
    notes.append(["renewal conversation scheduled for March"])

    path = tmp_path / "book.xlsx"
    wb.save(str(path))
    return path


def _write_image(path: Path, size: tuple[int, int] = (24, 24)) -> Path:
    """Save a real raster in the format implied by `path`'s suffix."""
    from PIL import Image

    if path.suffix.lower() in (".heic", ".heif"):
        import pillow_heif

        pillow_heif.register_heif_opener()
    Image.new("RGB", size, (255, 255, 255)).save(str(path))
    return path


def _ocr_stub(text: str = "ocr text") -> MagicMock:
    """Stub only the tesseract call — the binary isn't guaranteed anywhere."""
    stub = MagicMock()
    stub.image_to_string.return_value = text
    return stub


# --------------------------------------------------------------------------
# pptx
# --------------------------------------------------------------------------
async def test_pptx_extracts_slide_text_table_and_speaker_notes(deck: Path) -> None:
    result = await LocalExtractor().extract(deck, PPTX_MIME)

    assert "Quarterly Revenue" in result.text
    assert "Bookings grew eighteen percent" in result.text
    # Speaker notes are where a deck's actual argument usually lives.
    assert "Mention the Frankfurt outage" in result.text
    # Table cells are shape text too and must not be dropped.
    assert "Account" in result.text and "Acme" in result.text
    assert result.backend == "local"


async def test_pptx_routes_on_mime_when_the_suffix_lies(deck: Path, tmp_path: Path) -> None:
    """Mime is checked alongside suffix — a temp file may be named anything."""
    disguised = tmp_path / "upload-8f21.bin"
    disguised.write_bytes(deck.read_bytes())

    result = await LocalExtractor().extract(disguised, PPTX_MIME)

    assert "Quarterly Revenue" in result.text


async def test_pptx_slide_cap_truncates_and_says_so(deck: Path, monkeypatch) -> None:
    monkeypatch.setattr(local_mod, "_MAX_PPTX_SLIDES", 1)

    result = await LocalExtractor().extract(deck, PPTX_MIME)

    assert "Quarterly Revenue" in result.text  # slide 1 survived
    assert "Acme" not in result.text  # slide 2 was dropped
    assert "truncated" in result.text  # and the reader is told


async def test_pptx_refuses_a_file_over_the_office_ceiling(deck: Path, monkeypatch) -> None:
    """The ceiling is a pre-spend refusal, so it must RAISE, not return ''."""
    monkeypatch.setattr(local_mod, "_MAX_OFFICE_BYTES", 10)

    with pytest.raises(RuntimeError, match="pptx file too large"):
        await LocalExtractor().extract(deck, PPTX_MIME)


async def test_pptx_missing_python_pptx_raises_loudly(tmp_path: Path) -> None:
    path = tmp_path / "x.pptx"
    path.write_bytes(b"PK\x03\x04")

    with patch.dict(sys.modules, {"pptx": None}):
        with pytest.raises(RuntimeError, match="python-pptx not installed"):
            await LocalExtractor().extract(path, PPTX_MIME)


# --------------------------------------------------------------------------
# xlsx
# --------------------------------------------------------------------------
async def test_xlsx_extracts_sheet_names_and_cell_text(workbook: Path) -> None:
    result = await LocalExtractor().extract(workbook, XLSX_MIME)

    assert "Pipeline" in result.text  # sheet name
    assert "Notes" in result.text  # second sheet name
    assert "Acme" in result.text and "42000" in result.text
    assert "renewal conversation scheduled for March" in result.text


async def test_xlsx_routes_on_mime_when_the_suffix_lies(workbook: Path, tmp_path: Path) -> None:
    disguised = tmp_path / "upload-3c02.bin"
    disguised.write_bytes(workbook.read_bytes())

    result = await LocalExtractor().extract(disguised, XLSX_MIME)

    assert "Acme" in result.text


async def test_xlsx_per_sheet_cap_still_leaves_room_for_later_sheets(
    tmp_path: Path, monkeypatch
) -> None:
    """The load-bearing cap: one fat sheet must not starve the ones after it.

    A workbook's cost is not its file size — a generated sheet can hold a
    million rows inside a small zip. Without a PER-SHEET cap the first sheet
    consumes the whole budget and every later sheet is silently lost.
    """
    from openpyxl import Workbook

    wb = Workbook()
    fat = wb.active
    fat.title = "Fat"
    for i in range(400):
        fat.append([f"filler row {i} " + "x" * 60])
    tail = wb.create_sheet("Tail")
    tail.append(["the last sheet still matters"])
    path = tmp_path / "big.xlsx"
    wb.save(str(path))

    monkeypatch.setattr(local_mod, "_MAX_XLSX_CHARS_PER_SHEET", 500)
    result = await LocalExtractor().extract(path, XLSX_MIME)

    assert "truncated" in result.text
    assert "the last sheet still matters" in result.text


async def test_xlsx_sheet_count_cap_truncates_and_says_so(workbook: Path, monkeypatch) -> None:
    monkeypatch.setattr(local_mod, "_MAX_XLSX_SHEETS", 1)

    result = await LocalExtractor().extract(workbook, XLSX_MIME)

    assert "Acme" in result.text
    assert "renewal conversation scheduled for March" not in result.text
    assert "truncated" in result.text


async def test_xlsx_refuses_a_file_over_the_office_ceiling(workbook: Path, monkeypatch) -> None:
    monkeypatch.setattr(local_mod, "_MAX_OFFICE_BYTES", 10)

    with pytest.raises(RuntimeError, match="xlsx file too large"):
        await LocalExtractor().extract(workbook, XLSX_MIME)


async def test_xlsx_missing_openpyxl_raises_loudly(tmp_path: Path) -> None:
    path = tmp_path / "x.xlsx"
    path.write_bytes(b"PK\x03\x04")

    with patch.dict(sys.modules, {"openpyxl": None}):
        with pytest.raises(RuntimeError, match="openpyxl not installed"):
            await LocalExtractor().extract(path, XLSX_MIME)


# --------------------------------------------------------------------------
# html
# --------------------------------------------------------------------------
async def test_html_extracts_the_article_and_drops_the_chrome(tmp_path: Path) -> None:
    path = tmp_path / "post.html"
    path.write_text(ARTICLE_HTML, encoding="utf-8")

    result = await LocalExtractor().extract(path, "text/html")

    assert "A document that cannot be read" in result.text
    assert "Coverage is therefore not a nice to have" in result.text
    # Boilerplate is what trafilatura buys us over a tag strip.
    assert "home about contact login" not in result.text
    assert "all rights reserved" not in result.text


async def test_html_no_longer_indexes_raw_markup(tmp_path: Path) -> None:
    """Regression pin: before this branch, .html took the raw read_text path
    and the KB indexed the tags themselves."""
    path = tmp_path / "post.html"
    path.write_text(ARTICLE_HTML, encoding="utf-8")

    result = await LocalExtractor().extract(path, "text/html")

    assert "<p>" not in result.text
    assert "<article>" not in result.text
    assert "<html>" not in result.text


async def test_trafilatura_really_declines_the_fallback_fixtures() -> None:
    """Guard the guard: the two tests below only mean something if trafilatura
    actually returns None for these inputs.

    It does NOT for the obvious candidates — a bare `<p>hi</p>` document, or
    one wrapped in `<html><body>`, both extract cleanly. Fixtures like those
    make the fallback tests pass while never executing a line of the fallback,
    which is how two mutations to that code escaped this suite on the first
    sweep. A bodyless FRAGMENT is what trafilatura declines. If a future
    version starts accepting them, this fails first and names the reason
    instead of leaving the fallback quietly uncovered.
    """
    import trafilatura

    for fixture in (FRAGMENT_HTML, FRAGMENT_WITH_SCRIPT_HTML):
        assert (
            trafilatura.extract(
                fixture, output_format="markdown", include_links=True, include_tables=True
            )
            is None
        )


async def test_html_falls_back_to_tag_stripping_when_trafilatura_declines(
    tmp_path: Path,
) -> None:
    """trafilatura returning None is a CONTENT outcome, not a fault.

    It declines bodyless fragments. The fallback must still produce readable
    text rather than raw markup — and must not raise.
    """
    path = tmp_path / "fragment.html"
    path.write_text(FRAGMENT_HTML, encoding="utf-8")

    result = await LocalExtractor().extract(path, "text/html")

    assert "visible" in result.text
    assert "fragment text" in result.text
    assert "<" not in result.text and ">" not in result.text


async def test_html_fallback_drops_script_and_style_bodies(tmp_path: Path) -> None:
    path = tmp_path / "fragment.html"
    path.write_text(FRAGMENT_WITH_SCRIPT_HTML, encoding="utf-8")

    result = await LocalExtractor().extract(path, "text/html")

    assert "visible" in result.text
    assert "leak" not in result.text
    assert "crimson" not in result.text


async def test_html_routes_on_mime_when_the_suffix_lies(tmp_path: Path) -> None:
    path = tmp_path / "page.txt"
    path.write_text(ARTICLE_HTML, encoding="utf-8")

    result = await LocalExtractor().extract(path, "text/html")

    assert "<p>" not in result.text
    assert "A document that cannot be read" in result.text


async def test_html_charset_in_the_mime_does_not_break_routing(tmp_path: Path) -> None:
    path = tmp_path / "page.txt"
    path.write_text(ARTICLE_HTML, encoding="utf-8")

    result = await LocalExtractor().extract(path, "text/html; charset=utf-8")

    assert "<p>" not in result.text


async def test_html_refuses_a_file_over_the_html_ceiling(tmp_path: Path, monkeypatch) -> None:
    path = tmp_path / "post.html"
    path.write_text(ARTICLE_HTML, encoding="utf-8")
    monkeypatch.setattr(local_mod, "_MAX_HTML_BYTES", 10)

    with pytest.raises(RuntimeError, match="html file too large"):
        await LocalExtractor().extract(path, "text/html")


async def test_html_missing_trafilatura_raises_loudly(tmp_path: Path) -> None:
    """A missing dep is a packaging fault — never a quiet downgrade to regex."""
    path = tmp_path / "post.html"
    path.write_text(ARTICLE_HTML, encoding="utf-8")

    with patch.dict(sys.modules, {"trafilatura": None}):
        with pytest.raises(RuntimeError, match="trafilatura not installed"):
            await LocalExtractor().extract(path, "text/html")


# --------------------------------------------------------------------------
# raster formats — real decode, stubbed OCR
# --------------------------------------------------------------------------
@pytest.mark.parametrize(
    ("suffix", "pil_format"),
    [(".webp", "WEBP"), (".tiff", "TIFF"), (".tif", "TIFF"), (".heic", "HEIF")],
)
async def test_new_raster_formats_reach_the_ocr_path(
    tmp_path: Path, suffix: str, pil_format: str
) -> None:
    """PIL decodes the REAL bytes; only the tesseract call is stubbed.

    Asserting the FORMAT of the image handed to the stub is what proves the
    fixture was genuinely decoded — a test that only asserts the return value
    would pass against a file Pillow could not open at all.
    """
    path = _write_image(tmp_path / f"shot{suffix}")
    stub = _ocr_stub("text lifted from the image")

    with patch.dict(sys.modules, {"pytesseract": stub}):
        result = await LocalExtractor().extract(path, "")

    assert result.text == "text lifted from the image"
    (image,), _ = stub.image_to_string.call_args
    assert image.format == pil_format


async def test_heic_registers_the_pillow_heif_opener(tmp_path: Path) -> None:
    """Pillow cannot decode HEIF natively, so the opener must be registered.

    Asserted by call, not just by outcome: `register_heif_opener` is
    process-global, so once ANY earlier test has called it a decode would keep
    working even if this branch stopped registering. Only the call assertion
    fails when the registration is removed.
    """
    import pillow_heif

    path = _write_image(tmp_path / "photo.heic")
    stub = _ocr_stub()

    with patch.dict(sys.modules, {"pytesseract": stub}):
        with patch.object(
            pillow_heif, "register_heif_opener", wraps=pillow_heif.register_heif_opener
        ) as register:
            await LocalExtractor().extract(path, "")

    register.assert_called_once()


async def test_heic_missing_pillow_heif_raises_loudly(tmp_path: Path) -> None:
    path = tmp_path / "photo.heic"
    path.write_bytes(b"\x00\x00\x00\x18ftypheic")

    with patch.dict(sys.modules, {"pytesseract": _ocr_stub(), "pillow_heif": None}):
        with pytest.raises(RuntimeError, match="pillow-heif not installed"):
            await LocalExtractor().extract(path, "")


async def test_png_is_not_charged_the_heif_dependency(tmp_path: Path) -> None:
    """A png must keep working on a box with no pillow-heif installed."""
    path = _write_image(tmp_path / "shot.png")
    stub = _ocr_stub("still fine")

    with patch.dict(sys.modules, {"pytesseract": stub, "pillow_heif": None}):
        result = await LocalExtractor().extract(path, "")

    assert result.text == "still fine"


async def test_image_refuses_a_file_over_the_image_ceiling(tmp_path: Path, monkeypatch) -> None:
    path = _write_image(tmp_path / "shot.tiff")
    monkeypatch.setattr(local_mod, "_MAX_IMAGE_BYTES", 10)

    with patch.dict(sys.modules, {"pytesseract": _ocr_stub()}):
        with pytest.raises(RuntimeError, match="image file too large"):
            await LocalExtractor().extract(path, "")


# --------------------------------------------------------------------------
# ceilings, budget, and the explicit non-goals
# --------------------------------------------------------------------------
def test_the_declared_ceilings_are_finite_and_sane() -> None:
    """The per-test monkeypatches would happily pass against a ceiling of
    zero or of 2**64, so the shipped VALUES need their own assertion."""
    assert 1 * 1024 * 1024 <= local_mod._MAX_HTML_BYTES <= 64 * 1024 * 1024
    assert 1 * 1024 * 1024 <= local_mod._MAX_IMAGE_BYTES <= 128 * 1024 * 1024
    assert 1 * 1024 * 1024 <= local_mod._MAX_OFFICE_BYTES <= 256 * 1024 * 1024
    assert 1 <= local_mod._MAX_XLSX_SHEETS <= 1000
    assert 1_000 <= local_mod._MAX_XLSX_CHARS_PER_SHEET <= 1_000_000
    assert 1 <= local_mod._MAX_PPTX_SLIDES <= 10_000
    assert 10_000 <= local_mod._MAX_EXTRACTED_CHARS <= 5_000_000


def test_the_global_budget_truncates_with_a_visible_marker() -> None:
    text = local_mod._budgeted(["x" * 50], limit=10)
    content, _, marker = text.partition("\n")

    assert content == "x" * 10  # cut exactly at the limit, not near it
    assert "truncated" in marker


def test_the_global_budget_leaves_short_text_alone() -> None:
    assert local_mod._budgeted(["a", "b"], limit=100) == "a\nb"


@pytest.mark.parametrize("suffix", [".psd", ".cr2", ".nef", ".arw", ".ppt", ".xls"])
async def test_explicit_non_goals_keep_the_raw_read(tmp_path: Path, suffix: str) -> None:
    """psd and camera raw are out of scope by decision; .ppt/.xls are binary
    formats python-pptx and openpyxl cannot read. All must keep falling
    through to the raw text read rather than raising a parse error."""
    path = tmp_path / f"file{suffix}"
    path.write_text("plain bytes", encoding="utf-8")

    result = await LocalExtractor().extract(path, "")

    assert result.text == "plain bytes"
