# local.py — Local extraction adapter (pypdf + python-docx + pytesseract +
# python-pptx + openpyxl + trafilatura + Pillow/pillow-heif).
# Created: 2026-04-30 — Phase 1 of "Files as Knowledge" plan, Stage 1.A.
# Updated: 2026-08-29 (files-intelligence Track 1) — the dispatch table gained
#   .pptx (slides + speaker notes), .xlsx (sheet names + cell text), .html
#   (trafilatura), and the raster set .webp/.tif/.tiff/.heic/.heif routed into
#   the existing OCR path. Every new branch carries a size ceiling checked
#   BEFORE the library is handed the file. The old "line-for-line port of
#   _extract_file" framing no longer holds: the PDF/DOCX/image/VTT branches are
#   still that port untouched, but the module is now a superset of it.
"""LocalExtractor — wraps the traditional extraction libraries.

The PDF, DOCX, image and VTT branches remain a behavior-preserving port of the
old `_extract_file` in `agents/knowledge.py`; the rest of the table was added
later. `supports_mimes = {"*"}` keeps Local the chain's catch-all fallback.

Size ceilings — why they are byte checks, and why they run first
---------------------------------------------------------------
Extraction runs inside a bus listener (`uploads/listeners.py`) that catches
everything and continues, on a request-path thread. A pathological input
therefore cannot be allowed to consume the box; it has to be refused before any
library touches it. A timeout is the wrong instrument: by the time it fires the
memory is already allocated, and `openpyxl`/`python-pptx` do their damage in a
single un-interruptible C-level parse, so there is nothing to interrupt.

`path.stat().st_size` costs one syscall and runs before every `open`:

* ``_MAX_OFFICE_BYTES`` (50 MB) — OOXML is a zip. A deck or workbook that large
  is embedded media, not text; the text-heavy worst case (a 500-slide deck, a
  50-sheet financial model) sits comfortably under it. python-pptx and openpyxl
  both build an in-memory object graph several times the compressed size, so
  the multiplier, not the file, is what bounds us.
* ``_MAX_HTML_BYTES`` (10 MB) — trafilatura parses the whole DOM through lxml.
  10 MB of markup is ~100x a long article; past that it is a crawl dump.
* ``_MAX_IMAGE_BYTES`` (25 MB) — Pillow decodes to an uncompressed raster, so
  the compressed size understates the cost. Pillow's own decompression-bomb
  guard (`Image.MAX_IMAGE_PIXELS`) backstops the pixel dimension separately;
  this bounds the read.

Two different over-limit behaviours, deliberately:

* A **byte ceiling raises** `RuntimeError`. It is a pre-spend refusal — we
  cannot produce partial text from a file we declined to open, and a loud error
  naming the cap is the only thing that distinguishes "too big" from "empty" in
  the listener's log. Returning `""` here would be indistinguishable from a
  missing dependency, which is the exact failure this module has already
  shipped once (pypdf absent from the image: green build, every PDF silently
  empty, read as "the feature is off").
* **Structural caps truncate.** Sheet, slide and character limits are hit
  mid-parse, where partial text is strictly more useful than none, and the
  output carries a visible ``[truncated: ...]`` marker so a reader can tell
  a cap from a short document. ``_MAX_EXTRACTED_CHARS`` (100_000) matches the
  existing precedent in `pocketpaw/knowledge/ingest.py`.

The caps apply to the branches added in 2026-08; the ported PDF/DOCX/image
branches are left exactly as they were so their parity tests keep meaning what
they say.
"""

from __future__ import annotations

from pathlib import Path

from pocketpaw_ee.cloud.extraction.adapter import ExtractionResult

# --- Size ceilings (see the module docstring for the reasoning) -------------
_MAX_OFFICE_BYTES = 50 * 1024 * 1024
_MAX_HTML_BYTES = 10 * 1024 * 1024
_MAX_IMAGE_BYTES = 25 * 1024 * 1024

# Structural caps — hit mid-parse, so these truncate rather than raise.
_MAX_XLSX_SHEETS = 50
_MAX_XLSX_CHARS_PER_SHEET = 20_000
_MAX_PPTX_SLIDES = 500
_MAX_EXTRACTED_CHARS = 100_000

# Suffixes routed to the Pillow + pytesseract OCR path. png/jpg/jpeg are the
# original set; webp/tif/tiff decode natively on Pillow; heic/heif need the
# pillow-heif opener registered first. `.heic` is listed because that is what
# an iPhone actually writes — omitting it would leave the single most common
# phone-camera upload falling through to the raw-bytes read.
_IMAGE_SUFFIXES = (".png", ".jpg", ".jpeg", ".webp", ".tif", ".tiff", ".heic", ".heif")
_HEIF_SUFFIXES = (".heic", ".heif")

_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
_XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"


def _refuse_if_too_big(path: Path, ceiling: int, kind: str) -> None:
    """Raise before any library is handed the file. One stat(), no open().

    Loud on purpose: the caller (`uploads/listeners.py`) logs and continues, so
    a silent return would make an oversized file look identical to a missing
    dependency or an empty document.
    """
    size = path.stat().st_size
    if size > ceiling:
        raise RuntimeError(
            f"{kind} file too large for local extraction: "
            f"{size} bytes exceeds the {ceiling}-byte ceiling ({path.name})"
        )


def _budgeted(parts: list[str], limit: int = _MAX_EXTRACTED_CHARS) -> str:
    """Join collected text and truncate to the global budget with a marker."""
    text = "\n".join(parts)
    if len(text) <= limit:
        return text
    return text[:limit] + f"\n[truncated: extraction budget of {limit} characters reached]"


class LocalExtractor:
    """Wrapper around the offline extraction libraries.

    pypdf, python-docx, python-pptx, openpyxl, trafilatura and
    Pillow/pillow-heif + pytesseract. No network, so the chain can always
    fall back here.
    """

    name = "local"
    supports_mimes = {"*"}
    requires_network = False

    async def extract(self, path: Path, mime: str) -> ExtractionResult:
        text = await _extract_text(path, mime)
        return ExtractionResult(
            text=text,
            metadata={"path": str(path), "mime": mime},
            backend=self.name,
        )


async def _extract_text(path: Path, mime: str = "") -> str:
    """Extract text from PDF, DOCX, PPTX, XLSX, HTML, image or VTT.

    Anything unmatched falls back to a raw text read, as before.

    Routing precedence: explicit ``mime`` first, then file suffix. Mime
    wins because callers sometimes hand us a temp file whose extension
    doesn't match the source (e.g. a transcript stored under storage_key
    ``planner-XXXX.txt`` is really WebVTT — only the mime tells us to
    strip cue tags). Suffix is the fallback for legacy callers and for
    files where mime isn't reliably set (uploaded blobs). Every branch
    added in 2026-08 honours both, in that order.

    Deliberately NOT handled: ``.psd`` and camera raw. Both need heavy
    native decoders for output an OCR pass would mostly discard, and both
    are explicit non-goals of the Track 1 scope. They keep taking the raw
    read, which is what they did before.

    Also not handled: legacy ``.ppt`` / ``.xls``. python-pptx and openpyxl
    read the OOXML formats only — routing the binary ones here would trade
    a raw read for a confusing parse error.
    """
    file_path = str(path)
    suffix = path.suffix.lower()
    norm_mime = (mime or "").split(";", 1)[0].strip().lower()

    if suffix == ".pdf":
        try:
            from pypdf import PdfReader

            reader = PdfReader(file_path)
            return "\n".join(p.extract_text() or "" for p in reader.pages)
        except ImportError as exc:
            raise RuntimeError("pypdf not installed — run: pip install pypdf") from exc

    if suffix in (".docx", ".doc"):
        try:
            from docx import Document

            doc = Document(file_path)
            return "\n".join(p.text for p in doc.paragraphs)
        except ImportError as exc:
            raise RuntimeError("python-docx not installed — run: pip install python-docx") from exc

    if suffix == ".pptx" or norm_mime == _PPTX_MIME:
        _refuse_if_too_big(path, _MAX_OFFICE_BYTES, "pptx")
        # The import sits in its own try so an ImportError raised from INSIDE
        # python-pptx while parsing a real file can't be mis-reported as "the
        # library isn't installed" — the two need different fixes.
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise RuntimeError("python-pptx not installed — run: pip install python-pptx") from exc
        return _extract_pptx(Presentation(file_path))

    if suffix == ".xlsx" or norm_mime == _XLSX_MIME:
        _refuse_if_too_big(path, _MAX_OFFICE_BYTES, "xlsx")
        try:
            from openpyxl import load_workbook
        except ImportError as exc:
            raise RuntimeError("openpyxl not installed — run: pip install openpyxl") from exc
        # read_only streams rows instead of materialising the sheet; data_only
        # returns the cached value of a formula rather than "=SUM(A1:A9)",
        # which is the text a reader would actually search for.
        #
        # Handed a FILE OBJECT, not the path, and that is load-bearing:
        # load_workbook() validates the file EXTENSION and rejects anything
        # outside {.xlsx,.xlsm,.xltx,.xltm}. Uploads routinely arrive as a temp
        # file named `upload-8f21.bin` with the real type only in the mime, so
        # passing the path would make the mime route above dead code — it would
        # match, then raise InvalidFileException on a perfectly good workbook.
        # openpyxl skips the extension check for anything with a .read().
        with open(file_path, "rb") as handle:
            return _extract_xlsx(load_workbook(handle, read_only=True, data_only=True))

    if suffix in (".html", ".htm") or norm_mime == "text/html":
        _refuse_if_too_big(path, _MAX_HTML_BYTES, "html")
        return _extract_html(path.read_text(encoding="utf-8", errors="replace"))

    if suffix in _IMAGE_SUFFIXES:
        _refuse_if_too_big(path, _MAX_IMAGE_BYTES, "image")
        try:
            import pytesseract
            from PIL import Image
        except ImportError as exc:
            raise RuntimeError(
                "pytesseract not installed — run: pip install pytesseract Pillow"
            ) from exc

        if suffix in _HEIF_SUFFIXES:
            # Pillow does not decode HEIF natively (PIL.features.check("heif")
            # is False on Pillow 12.2). Without this the open() below raises
            # UnidentifiedImageError and an iPhone photo yields nothing.
            try:
                import pillow_heif
            except ImportError as exc:
                raise RuntimeError(
                    "pillow-heif not installed — run: pip install pillow-heif"
                ) from exc
            pillow_heif.register_heif_opener()

        return pytesseract.image_to_string(Image.open(file_path))

    if suffix == ".vtt" or norm_mime == "text/vtt":
        return _vtt_to_plain(path.read_text(encoding="utf-8", errors="replace"))

    return path.read_text(encoding="utf-8", errors="replace")


def _extract_pptx(presentation) -> str:
    """Slide text plus speaker notes, in slide order.

    Speaker notes are included because a deck's argument usually lives there —
    the slide carries three words and a chart, the notes carry the sentence
    that explains them. `has_notes_slide` is checked rather than reading
    `notes_slide` directly: python-pptx CREATES an empty notes slide as a side
    effect of the property access, which would mutate the file being read.
    """
    parts: list[str] = []
    truncated_slides = False

    for index, slide in enumerate(presentation.slides):
        if index >= _MAX_PPTX_SLIDES:
            truncated_slides = True
            break

        slide_parts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text_frame.text.strip()
                if text:
                    slide_parts.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        slide_parts.append("\t".join(cells))

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_parts.append(f"Notes: {notes}")

        if slide_parts:
            parts.append(f"# Slide {index + 1}")
            parts.extend(slide_parts)

    if truncated_slides:
        parts.append(f"[truncated: only the first {_MAX_PPTX_SLIDES} slides were read]")
    return _budgeted(parts)


def _extract_xlsx(workbook) -> str:
    """Sheet names plus cell text, capped per sheet and overall.

    The per-sheet cap is the load-bearing one. A workbook's cost is not its
    file size — one generated sheet can hold a million populated rows while
    the zip stays small, so a byte ceiling alone would let a single sheet
    exhaust the budget and starve every sheet after it. Capping per sheet
    means 40 sheets each contribute something.

    The workbook is closed in a finally: read_only mode holds the zip open.
    """
    parts: list[str] = []
    truncated_sheets = False
    try:
        for index, name in enumerate(workbook.sheetnames):
            if index >= _MAX_XLSX_SHEETS:
                truncated_sheets = True
                break

            parts.append(f"# {name}")
            used = 0
            for row in workbook[name].iter_rows(values_only=True):
                cells = [str(v).strip() for v in row if v is not None and str(v).strip()]
                if not cells:
                    continue
                line = "\t".join(cells)
                if used + len(line) > _MAX_XLSX_CHARS_PER_SHEET:
                    parts.append(
                        f"[truncated: sheet {name!r} exceeded "
                        f"{_MAX_XLSX_CHARS_PER_SHEET} characters]"
                    )
                    break
                parts.append(line)
                used += len(line)
    finally:
        workbook.close()

    if truncated_sheets:
        parts.append(f"[truncated: only the first {_MAX_XLSX_SHEETS} sheets were read]")
    return _budgeted(parts)


def _extract_html(html: str) -> str:
    """Article text via trafilatura, with a tag-strip fallback.

    Two distinct failure modes, handled differently on purpose:

    * trafilatura **missing** is a packaging fault — raise loudly, exactly as
      the pdf and docx branches do. Silently degrading to the regex stripper
      would hide a broken image behind slightly-worse output.
    * trafilatura **returning None** is normal. It is tuned for articles and
      declines boilerplate-only or very short documents. That is a content
      outcome, not a fault, so fall back to stripping tags (the idiom from
      `pocketpaw/knowledge/ingest.py`) rather than returning raw markup —
      before this branch existed, raw markup with all its tags is precisely
      what `.html` uploads were indexing into the KB.
    """
    import re

    try:
        import trafilatura
    except ImportError as exc:
        raise RuntimeError("trafilatura not installed — run: pip install trafilatura") from exc

    clean = (
        trafilatura.extract(html, output_format="markdown", include_links=True, include_tables=True)
        or ""
    )

    if not clean:
        clean = re.sub(r"<script[^>]*>.*?</script>", " ", html, flags=re.DOTALL | re.IGNORECASE)
        clean = re.sub(r"<style[^>]*>.*?</style>", " ", clean, flags=re.DOTALL | re.IGNORECASE)
        clean = re.sub(r"<[^>]+>", " ", clean)
        clean = re.sub(r"\s+", " ", clean).strip()

    return _budgeted([clean])


def _vtt_to_plain(vtt: str) -> str:
    """Strip a WebVTT blob down to readable speech for KB indexing.

    Keeps speaker-prefixed lines (``Speaker: text``) and drops the
    ``WEBVTT`` header, ``NOTE`` blocks, cue identifiers, and
    ``00:00:01.234 --> 00:00:05.678`` timestamp lines. Cue tags
    (``<v Speaker>...</v>``) are unwrapped into ``Speaker: ...``.
    Adjacent same-speaker turns are collapsed.

    The raw VTT remains the on-disk artifact for download; only the KB
    extraction sees the cleaned text. Embeddings + keyword search then
    score against speech instead of timestamps + markup noise.
    """
    import re

    cue_re = re.compile(r"<v\s+([^>]+)>([\s\S]*?)</v>", re.MULTILINE)
    timestamp_re = re.compile(r"^\s*\d{2}:\d{2}:\d{2}[.,]\d{3}\s*-->\s*")

    lines: list[str] = []
    last_speaker: str | None = None
    in_note = False
    for raw in vtt.splitlines():
        line = raw.strip()
        if not line:
            in_note = False
            continue
        if line == "WEBVTT" or line.startswith("WEBVTT "):
            continue
        if line.startswith("NOTE"):
            in_note = True
            continue
        if in_note:
            continue
        if timestamp_re.match(line):
            continue

        m = cue_re.search(line)
        if m:
            speaker = m.group(1).strip()
            text = re.sub(r"<[^>]+>", "", m.group(2)).strip()
            if not text:
                continue
            if speaker == last_speaker and lines:
                lines[-1] = f"{lines[-1]} {text}"
            else:
                lines.append(f"{speaker}: {text}")
                last_speaker = speaker
            continue

        # Plain line (no cue tag) — likely a single-speaker VTT. Keep it,
        # but drop bare cue identifiers (a single integer or short slug
        # on its own line, which VTT uses to label cues).
        if line.isdigit() or (len(line) <= 32 and "-" in line and " " not in line):
            continue
        lines.append(line)
        last_speaker = None

    return "\n".join(lines).strip()
